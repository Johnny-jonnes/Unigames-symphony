import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def setup_presentation():
    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0] # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = title_text
    subtitle.text = subtitle_text

def add_content_slide(prs, title_text, description, img_name):
    # Use blank layout to have full control
    slide_layout = prs.slide_layouts[5] # Title only
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title = slide.shapes.title
    title.text = title_text
    
    # Add description text box
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12.333)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = description
    p.font.size = Pt(18)
    
    # Add image
    img_path = os.path.join("docs", "img", img_name)
    if os.path.exists(img_path):
        # We want image width to be around 8 inches
        # Center it horizontally and put it below text
        pic_left = Inches((13.333 - 8) / 2)
        pic_top = Inches(2.5)
        slide.shapes.add_picture(img_path, pic_left, pic_top, width=Inches(8))
    else:
        print(f"Warning: Image {img_path} not found.")

def generate_manuel_pptx():
    prs = setup_presentation()
    add_title_slide(prs, "Manuel Utilisateur - UniGames", "Présentation des fonctionnalités clés de la plateforme")
    
    slides_data = [
        ("Connexion", "L'interface de connexion sécurisée permet aux utilisateurs d'accéder à la plateforme selon leur rôle.", "img_1.jpeg"),
        ("Tableau de bord", "Le Dashboard offre une vue d'ensemble avec des statistiques en temps réel pour l'édition sélectionnée.", "img_3.jpeg"),
        ("Gestion des Éditions", "Créez et configurez facilement de nouvelles éditions des Jeux Universitaires.", "img_4.jpeg"),
        ("Calendrier des Matchs", "Le calendrier des rencontres permet de programmer les matchs et de consulter les résultats.", "img_15.jpeg"),
        ("Programmation d'un Match", "Interface détaillée pour configurer la date, le lieu, et les équipes d'une rencontre.", "img_16.jpeg"),
        ("Classements", "Le tableau des classements est généré automatiquement en temps réel à partir des résultats de matchs.", "img_18.jpeg")
    ]
    
    for title, desc, img in slides_data:
        add_content_slide(prs, title, desc, img)
        
    prs.save(os.path.join("docs", "8_Manuel_Utilisateur_Presentation.pptx"))
    print("Généré : 8_Manuel_Utilisateur_Presentation.pptx")

def generate_maquette_pptx():
    prs = setup_presentation()
    add_title_slide(prs, "Maquettes UI/UX - UniGames", "Présentation du design et de l'expérience utilisateur")
    
    slides_data = [
        ("Design System & Couleurs", "Utilisation de couleurs professionnelles et modernes pour une interface claire et lisible.", "img_2.jpeg"),
        ("Composants Enterprise Cards", "Le Dashboard utilise des cartes avec statistiques clés pour mettre en évidence l'information importante.", "img_3.jpeg"),
        ("Formulaires Modernes", "Des formulaires interactifs et propres pour une saisie de données fluide.", "img_13.jpeg"),
        ("Listes et Tableaux de Données", "Des tableaux clairs avec badges de statut et actions intégrées.", "img_14.jpeg"),
        ("Visualisation des Classements", "Intégration de barres de progression pour illustrer visuellement le classement des équipes.", "img_18.jpeg")
    ]
    
    for title, desc, img in slides_data:
        add_content_slide(prs, title, desc, img)
        
    prs.save(os.path.join("docs", "9_Maquette_Presentation.pptx"))
    print("Généré : 9_Maquette_Presentation.pptx")

if __name__ == "__main__":
    generate_manuel_pptx()
    generate_maquette_pptx()
